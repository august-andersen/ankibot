"""File detection and text extraction for each supported file type."""

import base64
import csv
import io
from pathlib import Path

SUPPORTED_TEXT_EXTENSIONS = {".pdf", ".docx", ".pptx", ".csv", ".md"}
SUPPORTED_IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png"}
SUPPORTED_EXTENSIONS = SUPPORTED_TEXT_EXTENSIONS | SUPPORTED_IMAGE_EXTENSIONS

EXTENSION_LABELS = {
    ".pdf": "PDF",
    ".docx": "Word",
    ".pptx": "PowerPoint",
    ".csv": "CSV",
    ".md": "Markdown",
    ".jpg": "Image",
    ".jpeg": "Image",
    ".png": "Image",
}


def scan_directory(directory: Path) -> list[Path]:
    """Return supported files in the given directory (non-recursive), sorted by name."""
    files = []
    for p in sorted(directory.iterdir()):
        if p.is_file() and p.suffix.lower() in SUPPORTED_EXTENSIONS:
            files.append(p)
        elif p.is_file() and not p.name.startswith("."):
            print(f"  ⚠ Skipping unsupported file: {p.name}")
    return files


def is_image(path: Path) -> bool:
    return path.suffix.lower() in SUPPORTED_IMAGE_EXTENSIONS


def extract_pdf(path: Path) -> tuple[str, list[dict] | None]:
    """Extract text from PDF. Returns (text, image_pages) where image_pages
    is a list of base64-encoded page images if the PDF has no extractable text."""
    import fitz

    doc = fitz.open(path)
    text_parts = []
    for page in doc:
        text_parts.append(page.get_text())
    doc.close()

    full_text = "\n".join(text_parts).strip()
    if full_text:
        return full_text, None

    # Scanned PDF — render pages as images for vision
    doc = fitz.open(path)
    images = []
    for page in doc:
        pix = page.get_pixmap(dpi=150)
        img_bytes = pix.tobytes("png")
        images.append({
            "type": "image",
            "source": {
                "type": "base64",
                "media_type": "image/png",
                "data": base64.standard_b64encode(img_bytes).decode(),
            },
        })
    doc.close()
    return "", images


def extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    parts.append(" | ".join(cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Notes: {notes}]")
    return "\n".join(parts)


def extract_csv(path: Path) -> str:
    rows = []
    with open(path, newline="", encoding="utf-8") as f:
        reader = csv.reader(f)
        for row in reader:
            rows.append(" | ".join(row))
    return "\n".join(rows)


def extract_markdown(path: Path) -> str:
    return path.read_text(encoding="utf-8")


def extract_image(path: Path) -> dict:
    """Return a Claude API image content block for the given image file."""
    suffix = path.suffix.lower()
    media_types = {".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png"}
    data = base64.standard_b64encode(path.read_bytes()).decode()
    return {
        "type": "image",
        "source": {
            "type": "base64",
            "media_type": media_types[suffix],
            "data": data,
        },
    }


def extract_file(path: Path) -> tuple[str | None, list[dict] | None]:
    """Extract content from a file.
    Returns (text, image_blocks).
    - For text files: text is the content, image_blocks is None.
    - For images: text is None, image_blocks is a list with one image block.
    - For scanned PDFs: text is empty, image_blocks has page images.
    """
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            text, images = extract_pdf(path)
            if images:
                return None, images
            return text, None
        elif suffix == ".docx":
            return extract_docx(path), None
        elif suffix == ".pptx":
            return extract_pptx(path), None
        elif suffix == ".csv":
            return extract_csv(path), None
        elif suffix == ".md":
            return extract_markdown(path), None
        elif suffix in SUPPORTED_IMAGE_EXTENSIONS:
            return None, [extract_image(path)]
        else:
            return None, None
    except Exception as e:
        print(f"  ⚠ Failed to extract {path.name}: {e}")
        return None, None


def word_count(text: str) -> int:
    return len(text.split())
